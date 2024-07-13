from pptx import Presentation
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Define the path to the sample file
sample_file_path = './training/2024/2407/240721F.pptx'

# Function to extract and print text from all slides in a PowerPoint file
def print_slide_contents(ppt_file_path):
    # Load the presentation
    prs = Presentation(ppt_file_path)
    
    # Iterate through each slide in the presentation
    for slide_number, slide in enumerate(prs.slides, start=1):
        print(f"Slide {slide_number}:")
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Get the shape's type as a string
            # shape_type = MSO_SHAPE_TYPE.to_name(shape.shape_type)
            # Print the shape's name, type, and text if it has text
            shape_name = shape.name
            print(f"Name: {shape_name}")
            if hasattr(shape, "text"):
                print(f"Text: {shape.text}")
        print("-" * 20)  # Separator between slides

# Print the contents of the sample file
print_slide_contents(sample_file_path)