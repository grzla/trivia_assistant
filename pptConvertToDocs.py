from pptx import Presentation
import os

# Define the root directory path containing the PowerPoint files
root_directory_path = './training'
output_directory_path = os.path.join(root_directory_path, 'all')

# Ensure the output directory exists
os.makedirs(output_directory_path, exist_ok=True)

# Function to extract text from a slide, excluding 'Content Placeholder 14'
def extract_text_from_slide(slide):
    text_content = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.name != 'Content Placeholder 14' and shape.name != 'Text Placeholder 1':
            text_content.append(shape.text)
    return "\n".join(text_content)

# Function to find the game's date from the first slide or wherever it appears first
def find_games_date(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == 'Text Placeholder 1' and hasattr(shape, "text"):
                return shape.text
    return "Unknown Date"  # Return a default value if the date is not found

# Iterate through all files in the root directory
for root, dirs, files in os.walk(root_directory_path):
    for file in files:
        if file.endswith('.pptx'):
            file_path = os.path.join(root, file)
            prs = Presentation(file_path)
            
            # Find the game's date
            game_date = find_games_date(prs)
            
            # Extract text from slides 1, 3, and 5
            text_to_save = f"Game Date: {game_date}\n\n"  # Prepend the game's date
            for i in [0, 2, 4]:  # Slide numbers are 0-indexed in python-pptx
                if i < len(prs.slides):
                    text_to_save += extract_text_from_slide(prs.slides[i]) + "\n\n"
            
            # Define the output file path
            output_file_path = os.path.join(output_directory_path, os.path.splitext(file)[0] + '.txt')
            
            # Save the extracted text to a new text file
            with open(output_file_path, 'w') as text_file:
                text_file.write(text_to_save)

print("Conversion completed.")