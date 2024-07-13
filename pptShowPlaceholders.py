from pptx import Presentation
import os

# Define the path to the directory containing the sample files
file_path = './training'

# Function to extract and print text from specific placeholders in all PowerPoint files in a directory and its subdirectories
def print_specific_placeholders(directory_path):
	# Walk through the directory and its subdirectories
	for dirpath, dirnames, filenames in os.walk(directory_path):
		for file_name in filenames:
			# Check if the file is a PowerPoint file
			if file_name.endswith('.pptx'):
				full_path = os.path.join(dirpath, file_name)
				print(f"Processing file: {full_path}")
				# Load the presentation
				prs = Presentation(full_path)
				
				# Iterate through each slide in the presentation
				for slide_number, slide in enumerate(prs.slides, start=1):
					# Iterate through each shape in the slide
					for shape in slide.shapes:
						if shape.name in ['Content Placeholder 14', 'Text Placeholder 1']:
							if hasattr(shape, "text"):
								print(f"Slide {slide_number}, {shape.name}: {shape.text}")
				print("-" * 40)  # Separator between files

# Print the contents of the placeholders from all PowerPoint files in the directory and its subdirectories
print_specific_placeholders(file_path)